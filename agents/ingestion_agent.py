# agents/ingestion_agent.py

import os
import PyPDF2
from pptx import Presentation
import pandas as pd
from docx import Document
import markdown
from typing import List, Dict, Any

class IngestionAgent:
    """
    The IngestionAgent is responsible for parsing diverse document formats
    and preprocessing them into manageable text chunks.
    """
    def __init__(self):
        # Configuration for chunking. These can be tuned.
        self.chunk_size = 500
        self.chunk_overlap = 50

    def _read_pdf(self, file_path: str) -> str:
        """Reads text from a PDF file."""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page_num in range(len(reader.pages)):
                    text += reader.pages[page_num].extract_text() + "\n"
        except Exception as e:
            print(f"Error reading PDF {file_path}: {e}")
            return ""
        return text

    def _read_pptx(self, file_path: str) -> str:
        """Reads text from a PPTX file."""
        text = ""
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                text += f"--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error reading PPTX {file_path}: {e}")
            return ""
        return text

    def _read_csv(self, file_path: str) -> str:
        """Reads data from a CSV file and converts it to a string representation."""
        try:
            df = pd.read_csv(file_path)
            # Convert DataFrame to a string, useful for RAG
            text = df.to_string(index=False)
        except Exception as e:
            print(f"Error reading CSV {file_path}: {e}")
            return ""
        return text

    def _read_docx(self, file_path: str) -> str:
        """Reads text from a DOCX file."""
        text = ""
        try:
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        except Exception as e:
            print(f"Error reading DOCX {file_path}: {e}")
            return ""
        return text

    def _read_txt_md(self, file_path: str) -> str:
        """Reads text from a TXT or Markdown file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
                if file_path.endswith('.md'):
                    # Convert markdown to plain text for consistency
                    return markdown.markdown(content, strip_html=True)
                return content
        except Exception as e:
            print(f"Error reading TXT/MD {file_path}: {e}")
            return ""

    def _get_file_reader(self, file_extension: str):
        """Returns the appropriate reader function based on file extension."""
        if file_extension == '.pdf':
            return self._read_pdf
        elif file_extension == '.pptx':
            return self._read_pptx
        elif file_extension == '.csv':
            return self._read_csv
        elif file_extension == '.docx':
            return self._read_docx
        elif file_extension in ['.txt', '.md']:
            return self._read_txt_md
        else:
            return None

    def _split_text_into_chunks(self, text: str, source_file: str) -> List[Dict[str, Any]]:
        """Splits a long text into smaller, overlapping chunks."""
        chunks = []
        if not text:
            return chunks

        words = text.split()
        total_words = len(words)
        start_index = 0

        while start_index < total_words:
            end_index = min(start_index + self.chunk_size, total_words)
            chunk_words = words[start_index:end_index]
            chunk_content = " ".join(chunk_words)

            # Add metadata to each chunk
            chunks.append({
                "content": chunk_content,
                "source": source_file,
                "start_word_index": start_index,
                "end_word_index": end_index
            })

            # Move start index for next chunk, considering overlap
            start_index += (self.chunk_size - self.chunk_overlap)
            if start_index >= total_words: # Ensure we don't go past the end
                break
        return chunks

    def process_document(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Parses a document, extracts text, and splits it into chunks.
        Returns a list of dictionaries, where each dict represents a chunk
        with its content and metadata.
        """
        file_name = os.path.basename(file_path)
        file_extension = os.path.splitext(file_name)[1].lower()

        reader = self._get_file_reader(file_extension)
        if not reader:
            print(f"Unsupported file type: {file_extension}")
            return []

        print(f"Processing document: {file_name}")
        full_text = reader(file_path)
        if not full_text:
            print(f"Could not extract text from {file_name}")
            return []

        chunks = self._split_text_into_chunks(full_text, file_name)
        print(f"Extracted {len(chunks)} chunks from {file_name}")
        return chunks

# Example usage (for testing)
if __name__ == "__main__":
    # Create dummy files for testing
    if not os.path.exists("../documents"):
        os.makedirs("../documents")

    with open("../documents/test.txt", "w") as f:
        f.write("This is a test text file. It contains some sample content for chunking. " * 20)
    with open("../documents/test.md", "w") as f:
        f.write("# Markdown Test\n\nThis is a **markdown** file with some *formatting*.\n\n- Item 1\n- Item 2")
    # For PDF, PPTX, DOCX, CSV, you'd need actual files or mock them more complexly.
    # For simplicity, we'll just test with TXT/MD for now.

    ingestion_agent = IngestionAgent()

    # Test TXT
    txt_chunks = ingestion_agent.process_document("../documents/test.txt")
    print(f"\nTXT Chunks (first 2): {txt_chunks[:2]}")
    print(f"Total TXT chunks: {len(txt_chunks)}")

    # Test Markdown
    md_chunks = ingestion_agent.process_document("../documents/test.md")
    print(f"\nMarkdown Chunks (first 2): {md_chunks[:2]}")
    print(f"Total Markdown chunks: {len(md_chunks)}")

    # Clean up dummy files
    # os.remove("../documents/test.txt")
    # os.remove("../documents/test.md")
